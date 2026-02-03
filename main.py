# importing main web framework (fastapi), and libraries for handling file uploads,
# sending proper HTTP error responses, and enabling CORS (allowing comms between frontend and backend)
from fastapi import FastAPI, Request, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field, validator
from enum import Enum
import re
from typing import Optional, List
import json
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer

# handle requests, file extraction, hashing and timestamps
import requests
import os
import PyPDF2
from pptx import Presentation
import docx
from io import BytesIO
import hashlib
from datetime import datetime

# adding limits to request size (to prevent overloads)
from slowapi import Limiter
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
from fastapi.responses import JSONResponse
from fastapi import Header, Depends
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles

# adding logging for rate limiting
import logging
from logging.handlers import RotatingFileHandler

# loading environment variables from .env file
from dotenv import load_dotenv
import os
load_dotenv()
API_KEY_ENABLED = os.getenv("API_KEY_ENABLED", "false").lower() == "true"
API_KEY = os.getenv("API_KEY")

# initializing file upload limits
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
MAX_DOCS = 100  # max number of documents to store in memory

ALLOWED_MIME_TYPES = {
    ".pdf": "application/pdf",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
}

# initializing rate limiter
limiter = Limiter(key_func=get_remote_address)

# setting up logging for rate limiting
os.makedirs("logs", exist_ok=True)
#configuring logger
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        RotatingFileHandler("logs/api.log", maxBytes=10*1024*1024, backupCount=5),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)
logger.info("Starting AI Study Helper API...")

# creating the FastAPI app 
app = FastAPI(
    title = "AI Study Helper",
    description = "An AI-powered chatbot to assist with study materials.",
    version = "1.0.0"
)

# Health check endpoint
# simple endpoint to verify API and Ollama status
@app.get("/health")
def health_check():
    """Check if API and Ollama are running."""
    ollama_status = "unreachable"
    try:
        response = requests.post('http://localhost:11434/api/tags', timeout=2)
        if response.status_code == 200:
            ollama_status = "online"
    except:
        pass
    
    return {
        "status": "healthy",
        "ollama": ollama_status,
        "documents_stored": len(documents_db),
        "max_documents": MAX_DOCS
    }

# adding rate limiting to prevent abuse
app.state.limiter = limiter
@app.exception_handler(RateLimitExceeded)
async def rate_limit_handler(request, exc):
    logger.warning(f"Rate limit exceeded for IP: {request.client.host}, Path: {request.url.path}")
    return JSONResponse(
        status_code=429,
        content={
            "error": "Rate limit exceeded.",
            "detail": "You have sent too many requests. Please try again later.",
            "retry_after": "60 seconds"
        }
    )


# configuring CORS - allowing any frontend to call the api without restrictions
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# mounting static folder
app.mount("/static", StaticFiles(directory="static"), name = "static")
# serve index.html at root
@app.get("/")
def serve_frontend():
    return FileResponse("static/index.html")


# loading the embedder, a place to store the uploaded docs
# and a directory for storing the uploaded files
embedder = None
documents_db = {}
UPLOAD_DIR = "uploads"   
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Pydantic models for request bodies (Chat requests, flashcards, questions, study materials, document metadata)
# these define the structure of the data we expect in requests and return in responses

# client MUST send a document_id and material_type

class MaterialType(str, Enum):
    summary = "summary"
    flashcards = "flashcards"
    questions = "questions"
class StudyMaterialRequest(BaseModel):
    document_id: str = Field(..., min_length=32, max_length=32)
    material_type: MaterialType  # e.g., "pdf", "pptx", "docx"
    topic : Optional[str] = Field(None, max_length=100)
    
    @validator("topic")
    def sanitize_topic(cls, v):
        if v:
            return re.sub(r"[^\w\s\-]", "", v)
        return v

class ChatRequest(BaseModel):
    document_id: str = Field(..., min_length=32, max_length=32)
    question: str = Field(..., min_length=3, max_length=500)
    
    @validator("question")
    def sanitize_question(cls, v):
        v = v.replace("```", "")
        v = v.replace("<","").replace(">", "")
        v = v.replace("Ignore previous instructions.", "")
        return v.strip()

class FlashCard(BaseModel):
    front: str
    back: str
    topic: str

class Question(BaseModel):
    question: str
    type: str  # e.g., "multiple_choice", "short_answer"
    options: Optional[List[str]] = None
    answer: str
    explanation: str

class StudyMaterial(BaseModel):
    document_id: str
    material_type: str
    content: str
    created_at: str

class DocumentInfo(BaseModel):
    id: str
    filename: str
    upload_date: str
    material_type: str
    page_count: str
    text_length: str
    topics: List[str]

# startup event to load the embedder model
@app.on_event("startup")
async def startup_event():
    global embedder
    print("Loading embedder model...")
    embedder = SentenceTransformer('all-MiniLM-L6-v2')
    print("Embedder model loaded.")

# helper functions - converting binary files into raw text (critical for embedding and indexing)
# Extract text from PDF files
def extract_text_from_pdf(file_bytes: bytes) -> str:
    pdf_reader = PyPDF2.PdfReader(BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

# Extract text from pptx files
def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Extract text from docx files
def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(BytesIO(file_bytes))
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Method 1: 
def parse_flashcards_by_separator(raw_cards, doc):
    """Parse flashcards when separated by '---' """
    cards = []
    
    for card_block in raw_cards:
        card_block =  card_block.strip()
        
        if not card_block or "FRONT:" not in card_block or "BACK:" not in card_block:
            continue
        
        front_pos = card_block.find("FRONT:")
        back_pos = card_block.find("BACK:")
        
        if front_pos == -1 or back_pos == -1:
            continue
        
        front = card_block[front_pos + 6:back_pos].strip()
        back = card_block[back_pos + 5:].strip()
        
        if "FRONT:" in back:
            back = back[:back.find("FRONT:")].strip()
        
        if front and back and len(front) > 0 and len(back) > 0:
            cards.append({
                "front": front,
                "back": back,
                "topic": doc["topics"][0] if doc["topics"] else "General"
            })
    return cards[:10]

# Method 2:
def parse_flashcards_by_lines(response, doc):
    cards = []
    lines = response.split('\n')

    current_front = None
    current_back = None
    
    for line in lines:
        line = line.strip()
        
        if line.startswith("FRONT:"):
            if current_front and current_back:
                cards.append({  
                    "front": current_front,
                    "back": current_back,
                    "topic": doc["topics"][0] if doc["topics"] else "General"
                })
            
            current_front = line.replace("FRONT:", "").strip()
            current_back = None
            
        elif line.startswith("BACK:"):
            current_back = line.replace("BACK:", "").strip()
        
        elif current_back is not None and line is not line.startswith(("FRONT:", "BACK:", "---")):
            current_back += " " + line
            
    if current_front and current_back:
        cards.append({
            "front": current_front,
            "back": current_back,
            "topic": doc["topics"][0] if doc["topics"] else "General"
        })
    
    return cards[:10]

# Method 3: 
def parse_flashcards_aggressive(response, doc):
    """Aggressive parsing - search for all FRONT: and BACK: occurrences"""
    cards = []
    
    # Find all positions of FRONT: and BACK:
    import re
    
    # Use regex to find all flashcard patterns
    pattern = r'FRONT:\s*([^\n]+(?:\n(?!BACK:)[^\n]*)*)\s*BACK:\s*([^\n]+(?:\n(?!FRONT:)[^\n]*)*)'
    matches = re.finditer(pattern, response, re.MULTILINE)
    
    for match in matches:
        front = match.group(1).strip()
        back = match.group(2).strip()
        
        if front and back and len(front) > 0 and len(back) > 0:
            cards.append({
                "front": front,
                "back": back,
                "topic": doc["topics"][0] if doc["topics"] else "General"
            })
    
    return cards[:10]

# Function to chunk text into smaller pieces
# LLMs and embedding models have token limits, so we need to split large texts

# This function splits text into chunks of a specified size (in characters) 
# prevents overly large embedding inputs
def chunk_text(text: str, chunk_size: int = 500) -> List[str]:
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0
    for word in words:
        current_chunk.append(word)
        current_length += len(word) + 1  # +1 for the space
        if current_length >= chunk_size:
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0
    if current_chunk:
        chunks.append(" ".join(current_chunk))
        
    return chunks

# Function to create a FAISS index from text chunks
# this converts texts to embeddings, to float32, normalizes the vectors, and creates the FAISS index and adds vectors to the index.
# this enables semantic similarity search (CORE RAG ARCHITECTURE)
def create_faiss_index(chunks: List[str]):
    embeddings = embedder.encode(chunks, convert_to_tensor=False)
    embeddings = np.array(embeddings).astype("float32")
    faiss.normalize_L2(embeddings)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    return index, embeddings

# Sends prompt to Ollama API and retrieves AI-generated response, handles timeouts and errors
# Allows us to generate summaries, flashcards, questions, and chat responses
def generate_with_ollama(prompt: str, model: str = "llama3.2") -> str:        
    
    try:
        response = requests.post(
            'http://localhost:11434/api/generate',
            json={
                'model': model,
                'prompt': prompt,
                'stream': False,
                'options':{
                    'temperature':0.7,
                    'num_predict':1000
                }
            },
            timeout=120
        )
        if response.status_code == 200:
            return response.json()['response']
        else:
            return "Error: Ollama unavailable."
    except Exception as e:
        return f"Error: {str(e)}"

# Function to extract topics from text using AI
# Analyzes doc, extracts 3-5 main topics, and returns them as a list
# Used for metadata, flashcards, and question generation
def extract_topics(text: str) -> List[str]:
    # """Extract main topics from text (temporary - no AI)."""
    # # Simple keyword extraction until Ollama is set up
    # keywords = ['biology', 'chemistry', 'physics', 'math', 'history', 
    #             'science', 'anatomy', 'cell', 'molecule', 'equation']
    
    # text_lower = text.lower()
    # found_topics = []
    
    # for keyword in keywords:
    #     if keyword in text_lower:
    #         found_topics.append(keyword.capitalize())
    
    # # Return first 3 found topics, or generic ones
    # if found_topics:
    #     return found_topics[:3]
    # else:
    #     return ["General Study Notes", "Education", "Learning"]


# Uncomment below to use AI-based topic extraction when Ollama is set up

    """Extract main topics from text using AI."""
    prompt = f"""Analyze this text and extract 3-5 main topics/subjects covered.
Return ONLY a comma-separated list of topics, nothing else.

Text:
{text[:2000]}

Topics:"""
    try:
        response = generate_with_ollama(prompt)
        topics = [t.strip() for t in response.split(',')]
        return topics[:5]
    except Exception as e:
        print(f"Ollama error: {e}")
        return ["General Study Notes", "Education", "Learning"]

def sanitize_for_llm(text:str, max_length:int=3000) -> str:
    """Sanitize text for LLM input by removing special characters and limiting length."""
    text = text.replace("```", "")
    text = re.sub(r"[<>]", "", text)
    text = re.sub(r"Ignore previous instructions.", "", text)
    return text[:max_length].strip()

# API key verification dependency (for future use)
def verify_api_key(x_api_key: str = Header(None)):
    """Verify the provided API key."""
    if API_KEY_ENABLED:
        if not x_api_key or x_api_key != API_KEY:
            raise HTTPException(
                status_code=401, 
                detail="Invalid or missing API key."
            )
    return True

# Root endpoint - basic info about the API
# confirms the API is running and provides version and feature list
@app.get("/")
def root():
    return {"message": "Welcome to the AI Study Helper API!",
            "version": "1.0.0",
            "features": [
                "Upload study materials (PDF, PPTX, DOCX)",
                "Generate summary",
                "Create flashcards",
                "Generate practice questions",
                "Chat with your study materials"
            ]
        }

# Endpoint to upload a document
# Validates file type, extracts text, chunks it, creates FAISS index, extracts topics, and stores metadata
# This is the DOCUMENT UPLOAD AND PROCESSING CORE
@app.post("/upload", dependencies=[Depends(verify_api_key)])
@limiter.limit("5/minute") # limit to 5 uploads per minute 
async def upload_document(request: Request, file: UploadFile = File(...)):
    """
    Upload a study material document (PDF, PPTX, DOCX).
    Returns document ID for further processing.
    """
    file_bytes = await file.read()
    logger.info(f"Upload attempt - Filename: {file.filename}, Size: {len(file_bytes)} bytes")
    
    if len (documents_db) >= MAX_DOCS:
        raise HTTPException(
            status_code=429,
            detail="Document storage limit reached. "
        )
        
    #validating file type
    allowed_extensions = ['.pdf', '.pptx', '.docx']
    file_ext = os.path.splitext(file.filename)[1].lower()
    
    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400, detail= f"Unsupported file type. Allowed types: {', '.join(allowed_extensions)}"
            )
    
    # try-catch statements to handle file reading and text extraction
    try: 
        
        # checking file size
        if len(file_bytes) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413,
                detail="File too large. Maximum allowed size is 10 MB."
            )
        
        #validate mime type
        if file.content_type and file.content_type != ALLOWED_MIME_TYPES[file_ext]:
            raise HTTPException(
                status_code=400, detail="Invalid MIME type."
            )
        
        # extracting text based on file type
        if file_ext == '.pdf':
            text = extract_text_from_pdf(file_bytes)
        elif file_ext == '.pptx':
            text = extract_text_from_pptx(file_bytes)
        elif file_ext == '.docx':
            text = extract_text_from_docx(file_bytes)
        
        if not text.strip():
            raise HTTPException(
                status_code=400, detail="No extractable text found in the document."
            )
        
        # normalizing and processing text
        text = re.sub(r"\s+", " ", text).strip()
        
        # creating a doc id, chunking text, creating faiss index, extracting topics
        doc_id = hashlib.md5(file_bytes).hexdigest()
        chunks = chunk_text(text, chunk_size=500)
        
        # blocking duplicate uploads
        if doc_id in documents_db:
            raise HTTPException(
                status_code=409, detail="Document already uploaded."
            )
        # empty chunks check
        if not chunks:
            raise HTTPException(
                status_code=400, detail="Document text could not be chunked properly."
            )
        
        index, embeddings = create_faiss_index(chunks)
        topics = extract_topics(text)
        
        # storing document info in the in-memory db
        documents_db[doc_id] = {
            "id": doc_id,
            "filename": file.filename,
            "upload_date": datetime.utcnow().isoformat(),
            "text": text,
            "chunks": chunks,
            "index": index,
            "embeddings": embeddings,
            "topics": topics,
            "page_count": len(chunks),
            "text_length": len(text),
        }

        logger.info(f"Document uploaded successfully - ID: {doc_id}, Filename: {file.filename}")

        return {"document_id": doc_id,
                "filename": file.filename,
                "pages_processed": len(chunks),
                "topics_found": topics,
                "message": "Document uploaded and processed successfully."
                }
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error processing document: {str(e)}"
        )

# Listing uploaded documents endpoint
# returns metadata for all uploaded documents
@app.get("/documents", dependencies=[Depends(verify_api_key)])
@limiter.limit("30/minute") # limit to 30 document list requests per minute
def list_documents(request: Request):
    """
    List all uploaded documents with their metadata.
    """
    docs = []
    for doc_id, doc in documents_db.items():
        docs.append({
            "id": doc_id,
            "filename": doc["filename"],
            "upload_date": doc["upload_date"],
            "page_count": doc["page_count"],
            "topics": doc["topics"],
        })
    return {"documents": docs, "total": len(docs)}

# Endpoint to generate a summary for a document
# takes document id, retrieves text, constructs prompt, calls LLM, and returns summary
@app.post("/generate/summary", dependencies=[Depends(verify_api_key)])
@limiter.limit("10/minute") # limit to 10 summary requests per minute
async def generate_summary(request: Request, req: StudyMaterialRequest):
    """ Generate a summary for the specified document."""
    
    logger.info(f"Summary generation request for Document ID: {req.document_id}")
    
    if req.document_id not in documents_db:
        logger.warning(f"Document not found for summary generation - ID: {req.document_id}")
        raise HTTPException(status_code=404, detail="Document not found.")
    
    doc = documents_db[req.document_id]
    text = sanitize_for_llm(doc["text"], max_length=3000)  # limit to first 3000 chars for prompt size
    
    prompt = f"""Summarize the following study material in a concise manner. Include:
1. Main topics covered
2. Key concepts and definitions
3. Important facts, events, or formulas

Format as clear bullet points.
Content:
{text}
    
Summary:"""
    summary = generate_with_ollama(prompt)
    logger.info(f"Summary generated for Document ID: {req.document_id}")
    
    return {"document_id": req.document_id,
            "material_type": "summary",
            "content": summary,
            "created_at": datetime.now().isoformat()
            }

# Endpoint to generate flashcards for a document
# prompts LLM to create flashcards, parses response, and returns structured flashcard data, assigning topics automatically
@app.post("/generate/flashcards", dependencies=[Depends(verify_api_key)])
@limiter.limit("10/minute") # limit to 10 flashcard requests per minute
async def generate_flashcards(request: Request, req: StudyMaterialRequest):
    """ Generate flashcards for the specified document."""
    
    logger.info(f"Flashcard generation request for Document ID: {req.document_id}")
    
    if req.document_id not in documents_db:
        logger.warning(f"Document not found for flashcard generation - ID: {req.document_id}")
        raise HTTPException(status_code=404, detail="Document not found.")
    
    doc = documents_db[req.document_id]
    text = sanitize_for_llm(doc["text"], max_length=3000) # limit to first 3000 chars for prompt size
    
    prompt = f"""Create 10 flashcards based on the following study material.
Format EXACTLY as:
FRONT: question or term
BACK: answer or definition
---

FRONT: question or term
BACK: answer or definition
---

(repeat 10 times total)

Content to create flashcards from:
{text}

Flashcards:"""
    response = generate_with_ollama(prompt)
    
    logger.info(f"Ollama response length: {len(response)}, First 200 chars: {response[:200]}")
    
    # parsing the response into flashcard objects
    cards = []
    
    if '---' in response:
        raw_cards = response.split()
        cards = parse_flashcards_by_separator(raw_cards, doc)
        
    if len(cards) == 0:
        cards = parse_flashcards_by_lines(response, doc)
    
    if len(cards) == 0:
        cards = parse_flashcards_aggressive(response, doc)
        
    logger.info(f"Flashcards generated for Document ID: {req.document_id}")
    
    return {"document_id": req.document_id,
            "material_type": "flashcards",
            "flashcards": cards,
            "count": len(cards),
            "created_at": datetime.now().isoformat()
            }

# Endpoint to generate practice questions for a document
# requests LLM to create questions, parses response, and returns structured question data
@app.post("/generate/questions", dependencies=[Depends(verify_api_key)])
@limiter.limit("10/minute") # limit to 10 question requests per minute
async def generate_questions(request: Request, req: StudyMaterialRequest):
    """ Generate practice questions for the specified document."""
    
    logger.info(f"Question generation request for Document ID: {req.document_id}")
    
    if req.document_id not in documents_db:
        logger.warning(f"Document not found for question generation - ID: {req.document_id}")
        raise HTTPException(status_code=404, detail="Document not found.")
    
    doc = documents_db[req.document_id]
    text = sanitize_for_llm(doc["text"], max_length=3000)# limit to first 3000 chars for prompt size
    
    prompt = f"""Create 5 practice questions based on the following study material.
Format EXACTLY as:
Q: [question]
TYPE: [multiple_choice/short_answer]
OPTIONS: [A) option1, B) option2, C) option3, D) option4] (only for multiple_choice)
ANSWER: [correct answer]
EXPLANATION: [brief explanation of the answer, why it's correct]
---
(repeat for each question)

Content:
{text}

Questions:"""
    response = generate_with_ollama(prompt)
    
    # parsing the response into question objects
    questions = []
    q_texts = response.split('---')

    for block in q_texts:
        lines = [line.strip() for line in block.strip().split('\n') if line.strip()]
        
        q = {}
        for line in lines:
            if line.startswith("Q:"):
                q["question"] = line.replace("Q:", "").strip()
            elif line.startswith("TYPE:"):
                q["type"] = line.replace("TYPE:", "").strip().lower()
            elif line.startswith("OPTIONS:"):
                options_str = line.replace("OPTIONS:", "").strip()
                if options_str:
                    q["options"] = [opt.strip() for opt in options_str.split(",")]
            elif line.startswith("ANSWER:"):
                q["answer"] = line.replace("ANSWER:", "").strip()
            elif line.startswith("EXPLANATION:"):
                q["explanation"] = line.replace("EXPLANATION:", "").strip()
                
        if q.get("type") != "multiple_choice":
            q.pop("options", None)  # remove options if not multiple choice

        if "question" in q and "type" in q and "answer" in q:
            q.setdefault("type", "short_answer")
            q.setdefault("explanation", "")
            questions.append(q)
            
    questions = questions[:5]  # limit to first 5 questions
    
    logger.info(f"Questions generated for Document ID: {req.document_id}")
    return {"document_id": req.document_id,
            "material_type": "questions",
            "questions": questions,
            "count": len(questions),
            "created_at": datetime.now().isoformat()
            }

# Chat endpoint to interact with uploaded documents (RAG)
# Embedding-based retrieval of relevant chunks, constructs chat prompt, calls LLM, and returns response
@app.post("/chat", dependencies=[Depends(verify_api_key)])
@limiter.limit("20/minute") # limit to 20 chat requests per minute
async def chat_with_document(request: Request, req: ChatRequest):
    """ Ask questions about the uploaded document."""
    logger.info(f"Chat request for Document ID: {req.document_id}, Question: {req.question[:50]}")
    
    if req.document_id not in documents_db:
        logger.warning(f"Document not found for chat - ID: {req.document_id}")
        raise HTTPException(status_code=404, detail="Document not found.")
    
    doc = documents_db[req.document_id]
    
    # search for relevant chunks using FAISS
    query_vector = embedder.encode([req.question], convert_to_tensor=False)
    query_vector = np.array(query_vector).astype("float32")
    faiss.normalize_L2(query_vector)
    
    scores, indices = doc["index"].search(query_vector, k=3)  # retrieve top 3 relevant chunks
    
    # get relevant text chunks/context
    raw_context = "\n\n".join([doc["chunks"][idx] for idx in indices[0]])
    context = sanitize_for_llm(raw_context, max_length=1500)  # limit context size
    
    # generate an answer using the context
    prompt = f"""Use the following context from the study material to answer the question.
Context:
{context}

Question: {req.question}
Provide a clear, educational answer. If the context does not contain the answer, say so.

Answer:"""
    answer = generate_with_ollama(prompt)
    
    logger.info(f"Chat response generated for Document ID: {req.document_id}, Confidence Score: {float(scores[0][0])}")
    return {"document_id": req.document_id,
            "question": req.question,
            "answer": answer,
            "confidence_score": float(scores[0][0]),
            "sources_used": len(indices[0])
            }

# Endpoint to delete an uploaded document
# removes document and its data from in-memory db
@app.delete("/documents/{document_id}", dependencies=[Depends(verify_api_key)])
@limiter.limit("10/minute") # limit to 10 delete requests per minute
def delete_document(request: Request, document_id: str):
    """ Delete an uploaded document and its data."""
    logger.info(f"Delete request for Document ID: {document_id}")
    
    if document_id not in documents_db:
        logger.warning(f"Document not found for deletion - ID: {document_id}")
        raise HTTPException(status_code=404, detail="Document not found.")
    
    filename = documents_db[document_id]["filename"]
    del documents_db[document_id]
    
    logger.info(f"Document deleted successfully - ID: {document_id}, Filename: {filename}")
    return {"message": f"Document '{filename}' and its data have been deleted."}

# To run the app, use the command:
# uvicorn api:app --host
# This will start the FastAPI server on the specified host and port.
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
